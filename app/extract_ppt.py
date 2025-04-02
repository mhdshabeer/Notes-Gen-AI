from pptx import Presentation
import os
from PIL import Image
from io import BytesIO

def extract_ppt_content(ppt_path, output_folder="extracted_images"):
    """Extracts text and images from a PowerPoint file."""
    prs = Presentation(ppt_path)
    text_data = []

    # Create output folder for images if not exists
    os.makedirs(output_folder, exist_ok=True)

    for slide_idx, slide in enumerate(prs.slides):
        slide_text = f"Slide {slide_idx + 1}:\n"

        for shape_idx, shape in enumerate(slide.shapes):
            # Extract text
            if hasattr(shape, "text") and shape.text.strip():
                slide_text += f"{shape.text.strip()}\n"

            # Extract images
            if hasattr(shape, "image"):
                img = shape.image
                img_bytes = img.blob
                img_ext = img.ext or "png"
                img_path = os.path.join(output_folder, f"slide_{slide_idx+1}_img_{shape_idx+1}.{img_ext}")

                with open(img_path, "wb") as f:
                    f.write(img_bytes)

        text_data.append(slide_text.strip())

    return "\n\n".join(text_data)

# Example usage
if __name__ == "__main__":
    ppt_file = "example.pptx"  # Replace with your actual file path
    extracted_text = extract_ppt_content(ppt_file)
    print(extracted_text)
