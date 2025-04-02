# src/extract_text.py

import os
from pptx import Presentation
import PyPDF2
import pdfplumber

def extract_text_from_pptx(file_path, output_dir):
    """
    Extracts text from each slide in a PowerPoint file and saves it as individual text files.
    """
    presentation = Presentation(file_path)
    os.makedirs(output_dir, exist_ok=True)
    
    for i, slide in enumerate(presentation.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        
        text_file_path = os.path.join(output_dir, f"ppt_slide_{i}.txt")
        with open(text_file_path, "w", encoding="utf-8") as text_file:
            text_file.write("\n".join(slide_text))
        
        print(f"Extracted text from slide {i} to {text_file_path}")

def extract_text_from_pdf(file_path, output_dir):
    """
    Extracts text from each page in a PDF file and saves it as individual text files.
    """
    os.makedirs(output_dir, exist_ok=True)
    
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            page_text = page.extract_text()
            
            text_file_path = os.path.join(output_dir, f"pdf_page_{i}.txt")
            with open(text_file_path, "w", encoding="utf-8") as text_file:
                text_file.write(page_text or "No text found on this page.")
            
            print(f"Extracted text from page {i} to {text_file_path}")

def extract_text(file_path, output_dir):
    """
    Determines the file type and calls the appropriate text extraction function.
    """
    if file_path.endswith(".pptx"):
        extract_text_from_pptx(file_path, output_dir)
    elif file_path.endswith(".pdf"):
        extract_text_from_pdf(file_path, output_dir)
    else:
        raise ValueError("Unsupported file format. Only .pptx and .pdf are supported.")

if __name__ == "__main__":
    # Example usage
    input_file = "../input_files/example.pptx"
    output_directory = "../extracted_data/text/"
    extract_text(input_file, output_directory)