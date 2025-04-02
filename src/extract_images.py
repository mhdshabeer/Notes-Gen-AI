# src/extract_images.py

import os
from pptx import Presentation
from PIL import Image
from io import BytesIO
from pdf2image import convert_from_path

def extract_images_from_pptx(file_path, output_dir):
    """
    Extracts images from each slide in a PowerPoint file and saves them as individual image files.
    """
    presentation = Presentation(file_path)
    os.makedirs(output_dir, exist_ok=True)
    
    for slide_number, slide in enumerate(presentation.slides, start=1):
        image_count = 0
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Shape type 13 corresponds to images
                image = shape.image
                image_bytes = image.blob
                pil_image = Image.open(BytesIO(image_bytes))
                
                image_count += 1
                image_file_path = os.path.join(output_dir, f"ppt_slide_{slide_number}_image_{image_count}.png")
                pil_image.save(image_file_path)
                
                print(f"Extracted image from slide {slide_number} to {image_file_path}")

def extract_images_from_pdf(file_path, output_dir):
    """
    Extracts images from each page in a PDF file and saves them as individual image files.
    """
    os.makedirs(output_dir, exist_ok=True)
    
    pages = convert_from_path(file_path, dpi=300)
    for page_number, page in enumerate(pages, start=1):
        image_file_path = os.path.join(output_dir, f"pdf_page_{page_number}.png")
        page.save(image_file_path, "PNG")
        print(f"Extracted image from page {page_number} to {image_file_path}")

def extract_images(file_path, output_dir):
    """
    Determines the file type and calls the appropriate image extraction function.
    """
    if file_path.endswith(".pptx"):
        extract_images_from_pptx(file_path, output_dir)
    elif file_path.endswith(".pdf"):
        extract_images_from_pdf(file_path, output_dir)
    else:
        raise ValueError("Unsupported file format. Only .pptx and .pdf are supported.")

if __name__ == "__main__":
    # Example usage
    input_file = "../input_files/example.pptx"
    output_directory = "../extracted_data/images/"
    extract_images(input_file, output_directory)